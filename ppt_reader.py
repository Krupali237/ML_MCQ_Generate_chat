from typing import Tuple

from pptx import Presentation


def extract_text_from_pptx(path: str) -> Tuple[str, int]:
    """
    Read a PPTX file and return (text, num_slides).

    Raises ValueError for empty documents.
    """
    prs = Presentation(path)
    num_slides = len(prs.slides)

    if num_slides == 0:
        raise ValueError("The PowerPoint has no slides.")

    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)

    full_text = "\n".join(texts).strip()

    if not full_text:
        raise ValueError("Could not extract any text from the PPTX.")

    return full_text, num_slides

