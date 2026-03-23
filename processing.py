import os
import re
import uuid
import PyPDF2
from pptx import Presentation
import chromadb
from chromadb.config import Settings
from llm import get_embedding

UPLOAD_FOLDER = os.path.join(os.path.dirname(__file__), 'uploads')
CHROMA_PATH = os.path.join(os.path.dirname(__file__), 'chroma_db')

# Ensure directories exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(CHROMA_PATH, exist_ok=True)

# Initialize Chroma client
chroma_client = chromadb.PersistentClient(path=CHROMA_PATH)

def clean_text(text):
    # Remove extra whitespace, newlines, noise symbols
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'[^\w\s.,!?:;()\-%\'"]', '', text)
    return text.strip()

def extract_text_pdf(file_path):
    text = ""
    try:
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except Exception as e:
        print(f"Error reading PDF: {e}")
    return text

def extract_text_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading PPTX: {e}")
    return text

def chunk_text(text, max_words=150):
    sentences = re.split(r'(?<=[.!?]) +', text)
    chunks = []
    current_chunk = []
    current_length = 0
    
    for sentence in sentences:
        words = sentence.split()
        if current_length + len(words) > max_words and current_chunk:
            chunks.append(" ".join(current_chunk))
            current_chunk = [sentence]
            current_length = len(words)
        else:
            current_chunk.append(sentence)
            current_length += len(words)
            
    if current_chunk:
        chunks.append(" ".join(current_chunk))
        
    return chunks

def process_document(file_path):
    # Extract text based on extension
    ext = file_path.lower().split('.')[-1]
    if ext == 'pdf':
        text = extract_text_pdf(file_path)
    elif ext == 'pptx':
        text = extract_text_pptx(file_path)
    else:
        return None, "Unsupported file format"
    
    if not text.strip():
        return None, "No text found in document"
        
    cleaned_text = clean_text(text)
    chunks = chunk_text(cleaned_text)
    
    # Store in ChromaDB
    doc_id = str(uuid.uuid4())
    collection_name = f"doc_{doc_id.replace('-', '')}"
    
    try:
        # Create a new collection for this documenting
        collection = chroma_client.create_collection(name=collection_name)
        
        # Prepare data
        ids = [f"chunk_{i}" for i in range(len(chunks))]
        embeddings = []
        for chunk in chunks:
            # simple local retries or ignore failure
            emb = get_embedding(chunk)
            # If embedding fails, pad with zeros just to not crash, though not ideal
            if not emb:
               emb = [0] * 768 # nomic-embed-text size
            embeddings.append(emb)
            
        collection.add(
            ids=ids,
            embeddings=embeddings,
            documents=chunks
        )
        return doc_id, None
    except Exception as e:
        print(f"ChromaDB Error: {e}")
        return None, "Error storing data"

def retrieve_context(doc_id, query, top_k=10):
    collection_name = f"doc_{doc_id.replace('-', '')}"
    try:
        collection = chroma_client.get_collection(name=collection_name)
        query_embedding = get_embedding(query)
        
        results = collection.query(
            query_embeddings=[query_embedding],
            n_results=top_k
        )
        
        if results and 'documents' in results and results['documents']:
            return " ".join(results['documents'][0])
        return ""
    except Exception as e:
        print(f"Retrieval error: {e}")
        return ""

def get_all_chunks(doc_id, limit=200):
    collection_name = f"doc_{doc_id.replace('-', '')}"
    try:
        collection = chroma_client.get_collection(name=collection_name)
        res = collection.get(limit=limit)
        return res.get('documents', [])
    except Exception as e:
        print(f"Error fetching chunks: {e}")
        return []
